"""import collections
import collections.abc
from pptx import Presentation
from PyPDF2 import PdfReader
from rake_nltk import Rake

class Keywords:
    
    def pdf_text(self,path):
        reader = PdfReader(f"{path}")
        page = reader.pages[0]
        # print(page.extract_text())
        no_of_pages = len(reader.pages)
        text_list = ""
        for i in range(no_of_pages):
            page = reader.pages[i]
            text = page.extract_text()
            filter = ''.join([chr(i) for i in range(1, 32)])
            text_list += (text.translate(str.maketrans('', '', filter)))
        return text_list


    def pptx_text(self,path):
        file = open(f'{path}', "rb")
        prs = Presentation(file)
        text_list = []
        for slide in prs.slides:
            t = ""
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        t = t + run.text
            if t != "":
                text_list.append(t)
        return text_list

    def keyword_extract(self):
        #final_path = "C:\Stuffing\ImPing\COEP\Enigma_WC_Report.pdf"
        final_path="C:/Users/ameys/Downloads/cocppt.odp"
        l = len(final_path)
        if final_path[l-1] == 'x':
            output =self. pptx_text(final_path)
            print(output)
        else:
            output = self.pdf_text(final_path)
            print(output)

        rake = Rake()
        rake.extract_keywords_from_text(output)
        print(rake.get_ranked_phrases())
    # __containsi
o=Keywords()
o.keyword_extract()"""
"""import collections
from pptx import Presentation
from PyPDF2 import PdfFileReader
from rake_nltk import Rake

class Keywords:
    
    def pdf_text(self, path):
        reader = PdfFileReader(path)
        page = reader.getPage(0)
        # print(page.extract_text())
        no_of_pages = reader.getNumPages()
        text_list = ""
        for i in range(no_of_pages):
            page = reader.getPage(i)
            text = page.extractText()
            filter = ''.join([chr(i) for i in range(1, 32)])
            text_list += (text.translate(str.maketrans('', '', filter)))
        return text_list


    def pptx_text(self, path):
        file = open(path, "rb")
        prs = Presentation(file)
        text_list = []
        for slide in prs.slides:
            t = ""
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        t = t + run.text
            if t != "":
                text_list.append(t)
        return text_list

    def keyword_extract(self):
        final_path = "C:/Users/ameys/Downloads/cocppt.odp"
        l = len(final_path)
        if final_path[l-1] == 'x':
            output = self.pptx_text(final_path)
            print(output)
        else:
            output = self.pdf_text(final_path)
            print(output)

        rake = Rake()
        rake.extract_keywords_from_text(output)
        print(rake.get_ranked_phrases())

    def __contains__(self, item):
        return item in self.keywords

o = Keywords()
o.keyword_extract()
"""

# from pypdf import PdfReader
import PyPDF2
from pptx import Presentation
import glob
import nltk
from rake_nltk import Rake

text_list=[]

def read_pdf(path):
    reader = PyPDF2.PdfReader(path)
    page = reader.pages[0]
    no_of_pages = len(reader.pages)
    # text_list = []
    for i in range(no_of_pages):
        page = reader.pages[i]
        text = page.extract_text()
        filter = ''.join([chr(i) for i in range(1, 32)])
        #self.text_list.append(text.translate(str.maketrans('', '', filter)))
        text_list.append(text.translate(str.maketrans('', '', filter)))
    return text_list


def extract_keywords(fullText):
    print(text_list)
    r1 = Rake()
    keywords = []
    for text in fullText:
        print(text)
        r1.extract_keywords_from_text(text)
        keywords.append(r1.get_ranked_phrases_with_scores())
    return keywords

def extract_summary(fullText):
    summary = []
    for text in fullText:
        summary.append(nltk.sent_tokenize(text))
    return summary

def extract_text(path):
    if True:
        text_list = read_pdf(path)
        return extract_keywords(text_list)     
  
path = r"C:\Users\DELL\Desktop\Resume\Omkar_W_Resume.pdf"

print(extract_text(path))